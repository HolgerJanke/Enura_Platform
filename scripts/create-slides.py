from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

DARK = RGBColor(0x11, 0x18, 0x27)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x6B, 0x72, 0x80)
LIGHT_BG = RGBColor(0xF9, 0xFA, 0xFB)
PRIMARY = RGBColor(0x1A, 0x56, 0xDB)
ACCENT = RGBColor(0xF3, 0xA9, 0x17)
GREEN = RGBColor(0x05, 0x96, 0x69)
RED = RGBColor(0xDC, 0x26, 0x26)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)
BORDER = RGBColor(0xE5, 0xE7, 0xEB)
TEAL = RGBColor(0x0D, 0x94, 0x88)

def add_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_box(slide, left, top, w, h, fill_color, border_color=None, radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    if radius is not None:
        shape.adjustments[0] = radius
    return shape

def add_text(slide, left, top, w, h, text, size=14, color=DARK, bold=False, align=PP_ALIGN.LEFT, font_name='Arial'):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 1: Title
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_box(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), PRIMARY)

add_text(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(1),
         'ENURA PLATFORM', 16, GRAY)
add_text(slide, Inches(1.5), Inches(2.3), Inches(10), Inches(1.2),
         'Architektur-Overview', 44, WHITE, bold=True)
add_text(slide, Inches(1.5), Inches(3.5), Inches(10), Inches(0.6),
         'White-Label SaaS  |  Multi-Tenant  |  Branchenunabhangig  |  Bot-ready', 18, RGBColor(0x9C, 0xA3, 0xAF))

add_text(slide, Inches(1.5), Inches(4.8), Inches(10), Inches(1.0),
         'Eine Plattform — beliebig viele Firmen.\nJede Firma: eigenes Branding, eigene Module, eigene Connectors, eigene Bots.', 16, RGBColor(0x9C, 0xA3, 0xAF))

add_text(slide, Inches(1.5), Inches(6.3), Inches(10), Inches(0.5),
         'Stand: Mai 2026  |  Branch: feat/frontend-redesign', 13, GRAY)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 2: Konzept
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_box(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), PRIMARY)

add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.5),
         'Das Konzept: 1 Plattform, N Firmen', 28, DARK, bold=True)

add_text(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.5),
         'Jede Firma ist ein eigener Tenant mit eigener Subdomain. Alles wird pro Firma in der Datenbank konfiguriert.', 15, GRAY)

cols = [
    ('Firma A\nfirma-a.enura-group.com', PRIMARY,
     ['Blaues Branding', 'Vertrieb + Projekte aktiv', 'CRM-Connector angebunden', 'Setter-Bot aktiv']),
    ('Firma B\nfirma-b.enura-group.com', GREEN,
     ['Grunes Branding', 'Nur Vertrieb + Finance', 'Buchhaltungs-Connector', 'Finance-Bot aktiv']),
    ('Firma C\nfirma-c.enura-group.com', PURPLE,
     ['Lila Branding', 'Alle Module freigeschaltet', '3 Connectors aktiv', 'Alle Bots aktiv']),
]

for i, (title, color, features) in enumerate(cols):
    x = Inches(0.5 + i * 4.2)
    add_box(slide, x, Inches(1.9), Inches(3.8), Inches(4.0), LIGHT_BG, BORDER, 0.03)
    add_box(slide, x, Inches(1.9), Inches(3.8), Inches(0.08), color)
    lines = title.split('\n')
    add_text(slide, x + Inches(0.3), Inches(2.15), Inches(3.2), Inches(0.5),
             lines[0], 18, DARK, bold=True)
    add_text(slide, x + Inches(0.3), Inches(2.6), Inches(3.2), Inches(0.4),
             lines[1], 10, color, font_name='Consolas')
    for j, feat in enumerate(features):
        add_text(slide, x + Inches(0.3), Inches(3.2 + j * 0.55), Inches(3.2), Inches(0.4),
                 f'•  {feat}', 13, DARK)

add_box(slide, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.8), RGBColor(0xEF, 0xF6, 0xFF), BORDER, 0.02)
add_text(slide, Inches(0.8), Inches(6.3), Inches(11.5), Inches(0.6),
         'Gleicher Code, gleiche Infrastruktur — nur die Konfiguration (DB + Env Vars) bestimmt was jede Firma sieht.', 14, PRIMARY, bold=True)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 3: Tech Stack
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_box(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), PRIMARY)

add_text(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.5),
         'Tech Stack', 28, DARK, bold=True)

techs = [
    ('Next.js 14', 'App Router, SSR, Middleware\nServer Components + Client Hydration', PRIMARY),
    ('TypeScript', 'Monorepo: apps/web, apps/api\npackages/types, packages/ui', RGBColor(0x31, 0x78, 0xC6)),
    ('Supabase', 'PostgreSQL, Auth, RLS, Storage\nREST API + Realtime', GREEN),
    ('Tailwind CSS', 'CSS Variables = Brand Tokens\nPro Firma dynamisch via Middleware', PURPLE),
    ('Turborepo', 'Build-Pipeline, Caching\n@enura/types als shared Package', ACCENT),
    ('Vercel', 'Hosting, Edge Middleware\nPreview Deployments pro Branch', DARK),
]

for i, (title, desc, color) in enumerate(techs):
    col = i % 3
    row = i // 3
    x = Inches(0.8 + col * 4.0)
    y = Inches(1.4 + row * 2.7)
    add_box(slide, x, y, Inches(3.6), Inches(2.2), LIGHT_BG, BORDER, 0.03)
    add_box(slide, x, y, Inches(3.6), Inches(0.08), color)
    add_text(slide, x + Inches(0.3), y + Inches(0.3), Inches(3.0), Inches(0.5),
             title, 18, DARK, bold=True)
    add_text(slide, x + Inches(0.3), y + Inches(0.85), Inches(3.0), Inches(1.2),
             desc, 13, GRAY)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 4: Multi-Tenant Flow
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_box(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), PRIMARY)

add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.5),
         'Multi-Tenant: Wie erkennt die Plattform die Firma?', 28, DARK, bold=True)

boxes_flow = [
    ('Browser\n{firma}.enura-group.com', Inches(0.5), PRIMARY),
    ('Middleware\nSubdomain = Firma\nBranding aus DB', Inches(3.5), ACCENT),
    ('Supabase\ncompanies-Tabelle\ncompany_branding', Inches(6.5), GREEN),
    ('Page Render\nFirmen-Branding\nFirmen-Module', Inches(9.5), PURPLE),
]

for (label, x, color) in boxes_flow:
    box = add_box(slide, x, Inches(1.4), Inches(2.6), Inches(1.8), color, radius=0.04)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(label.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(13 if i == 0 else 11)
        p.font.color.rgb = WHITE
        p.font.bold = (i == 0)
        p.font.name = 'Arial'
        p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(12)

for x in [Inches(3.1), Inches(6.1), Inches(9.1)]:
    add_text(slide, x, Inches(2.0), Inches(0.4), Inches(0.5), '→', 24, GRAY, align=PP_ALIGN.CENTER)

add_text(slide, Inches(0.8), Inches(3.6), Inches(10), Inches(0.4),
         'Was ist pro Firma individuell konfiguriert?', 18, DARK, bold=True)

per_tenant = [
    ('Branding', 'Farben, Font, Logo, Border-Radius — alles in company_branding Tabelle', PRIMARY),
    ('Module', 'Welche Sidebar-Eintrage sichtbar — uber role_permissions Tabelle', GREEN),
    ('Connectors', 'Welche externen Systeme angebunden — uber company_connectors Tabelle', ACCENT),
    ('Bots', 'Welche KI-Bots laufen — pro Firma konfigurierbar (next step)', PURPLE),
    ('User & Rollen', 'Eigene Mitarbeiter, eigene Rollen-Zuweisung — profiles + profile_roles', TEAL),
]

for i, (title, desc, color) in enumerate(per_tenant):
    y = Inches(4.15 + i * 0.58)
    add_box(slide, Inches(0.8), y + Inches(0.05), Inches(0.12), Inches(0.35), color)
    add_text(slide, Inches(1.15), y, Inches(2.0), Inches(0.4), title, 14, DARK, bold=True)
    add_text(slide, Inches(3.2), y, Inches(9.0), Inches(0.4), desc, 13, GRAY)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 5: Rollen & Permissions
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_box(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), PRIMARY)

add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.5),
         'Rollen & Permissions', 28, DARK, bold=True)

add_text(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.5),
         'Jede Firma legt fest welche Rollen es gibt und welche Module freigeschaltet sind.', 14, GRAY)

# Roles
add_box(slide, Inches(0.5), Inches(1.6), Inches(3.5), Inches(5.3), LIGHT_BG, BORDER, 0.02)
add_text(slide, Inches(0.8), Inches(1.75), Inches(3), Inches(0.4), 'Standard-Rollen', 16, DARK, bold=True)

roles = [
    ('holding_admin', 'Plattform-Admin (Enura)'),
    ('super_user', 'Firmen-Admin'),
    ('geschaeftsfuehrung', 'GF / Management'),
    ('teamleiter', 'Teamleiter'),
    ('mitarbeiter', 'Mitarbeiter'),
    ('buchhaltung', 'Buchhaltung'),
]
for i, (key, label) in enumerate(roles):
    y = Inches(2.35 + i * 0.7)
    add_box(slide, Inches(0.8), y, Inches(2.9), Inches(0.55), WHITE, BORDER, 0.05)
    add_text(slide, Inches(1.0), y + Inches(0.02), Inches(2.5), Inches(0.3), key, 11, PRIMARY, bold=True, font_name='Consolas')
    add_text(slide, Inches(1.0), y + Inches(0.27), Inches(2.5), Inches(0.25), label, 10, GRAY)

# Permissions
add_box(slide, Inches(4.3), Inches(1.6), Inches(4.0), Inches(5.3), LIGHT_BG, BORDER, 0.02)
add_text(slide, Inches(4.6), Inches(1.75), Inches(3), Inches(0.4), 'Module (an/aus pro Firma)', 16, DARK, bold=True)

perms = [
    ('module:setter', 'Telefonakquise'),
    ('module:berater', 'Aussendienst'),
    ('module:leads', 'Lead-Verwaltung'),
    ('module:innendienst', 'Innendienst'),
    ('module:bau', 'Projekte & Montage'),
    ('module:finance', 'Finanzen'),
    ('module:reports', 'Analytics & Reports'),
    ('module:ai', 'KI-Features & Bots'),
    ('module:admin', 'Einstellungen'),
]
for i, (perm, label) in enumerate(perms):
    y = Inches(2.35 + i * 0.5)
    add_text(slide, Inches(4.6), y, Inches(2.0), Inches(0.4), perm, 11, DARK, font_name='Consolas')
    add_text(slide, Inches(6.7), y, Inches(1.5), Inches(0.4), label, 11, GRAY)

# How it works
add_box(slide, Inches(8.6), Inches(1.6), Inches(4.2), Inches(5.3), LIGHT_BG, BORDER, 0.02)
add_text(slide, Inches(8.9), Inches(1.75), Inches(3.5), Inches(0.4), 'Wie es funktioniert', 16, DARK, bold=True)

how_lines = [
    ('1.', 'User loggt ein → Session\n   mit Rollen + Permissions'),
    ('2.', 'Sidebar zeigt nur Module\n   die der User sehen darf'),
    ('3.', 'Seiten prufen Permission\n   bevor sie rendern'),
    ('4.', 'Firma A sieht anderes\n   Menu als Firma B'),
    ('5.', 'Plattform-Admin sieht\n   immer alles'),
]
for i, (num, text) in enumerate(how_lines):
    y = Inches(2.4 + i * 0.95)
    add_text(slide, Inches(8.9), y, Inches(0.4), Inches(0.4), num, 14, PRIMARY, bold=True)
    add_text(slide, Inches(9.3), y, Inches(3.3), Inches(0.8), text, 12, DARK)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 6: Prozesshaus
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_box(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), PRIMARY)

add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.5),
         'Prozesshaus: Der Modulbaukasten', 28, DARK, bold=True)

add_text(slide, Inches(0.8), Inches(0.95), Inches(10), Inches(0.4),
         'Jede Firma schaltet nur die Module frei die sie braucht.', 14, GRAY)

add_box(slide, Inches(1.0), Inches(1.5), Inches(7.0), Inches(0.7), DARK, radius=0.03)
add_text(slide, Inches(1.0), Inches(1.55), Inches(7.0), Inches(0.6),
         'ENURA PROZESSHAUS', 16, WHITE, bold=True, align=PP_ALIGN.CENTER)

procs = [
    ('P1', 'Vertrieb &\nAkquise', '/leads', PRIMARY),
    ('P2', 'Projekt-\nmanagement', '/projects', RGBColor(0x04, 0x7B, 0x57)),
    ('P3', 'Montage &\nTechnik', '/processes', RGBColor(0x92, 0x40, 0x0E)),
]
for i, (code, name, route, color) in enumerate(procs):
    x = Inches(1.0 + i * 2.35)
    add_box(slide, x, Inches(2.4), Inches(2.15), Inches(2.0), color, radius=0.03)
    add_text(slide, x, Inches(2.5), Inches(2.15), Inches(0.5), code, 22, WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x, Inches(3.0), Inches(2.15), Inches(0.7), name, 13, WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, x, Inches(3.7), Inches(2.15), Inches(0.4), route, 10, RGBColor(0xBF, 0xDB, 0xFE), align=PP_ALIGN.CENTER, font_name='Consolas')

sups = [
    ('S1', 'Analytics', '/analytics', PURPLE),
    ('S2', 'Finanzen &\nControlling', '/controlling', RGBColor(0xB4, 0x54, 0x09)),
    ('', 'Einstellungen', '/settings', GRAY),
]
for i, (code, name, route, color) in enumerate(sups):
    x = Inches(1.0 + i * 2.35)
    add_box(slide, x, Inches(4.6), Inches(2.15), Inches(2.0), color, radius=0.03)
    add_text(slide, x, Inches(4.7), Inches(2.15), Inches(0.5), code if code else '', 22, WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x, Inches(5.2), Inches(2.15), Inches(0.6), name, 13, WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, x, Inches(5.85), Inches(2.15), Inches(0.4), route, 10, RGBColor(0xBF, 0xDB, 0xFE), align=PP_ALIGN.CENTER, font_name='Consolas')

# Right: Beispiel
add_box(slide, Inches(8.3), Inches(1.5), Inches(4.5), Inches(5.3), LIGHT_BG, BORDER, 0.02)
add_text(slide, Inches(8.6), Inches(1.65), Inches(4.0), Inches(0.4), 'Beispiel', 16, DARK, bold=True)

examples = [
    ('Firma will nur Vertrieb + Finance?', DARK, False),
    ('→ Nur P1 + S2 aktivieren.\n   Rest verschwindet aus der Sidebar.', PRIMARY, False),
    ('', GRAY, False),
    ('Firma braucht alle Module?', DARK, False),
    ('→ Alle Permissions vergeben.\n   Voller Zugriff.', PRIMARY, False),
    ('', GRAY, False),
    ('Neue Firma onboarden?', DARK, False),
    ('→ Subdomain + Branding setzen\n→ Rollen + Module in DB\n→ Connectors konfigurieren\n→ Fertig', PRIMARY, False),
]

y_pos = Inches(2.2)
for text, color, _ in examples:
    if text == '':
        y_pos += Inches(0.2)
        continue
    is_arrow = text.startswith('→')
    add_text(slide, Inches(8.6), y_pos, Inches(4.0), Inches(0.85),
             text, 11 if is_arrow else 13, color, bold=not is_arrow)
    y_pos += Inches(0.85) if is_arrow else Inches(0.4)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 7: Connector System
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_box(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), PRIMARY)

add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.5),
         'Connector System: Beliebig erweiterbar', 28, DARK, bold=True)

add_text(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.5),
         'Jede Firma bindet ihre eigenen Systeme an. Neue Connectors = neuer Worker in apps/api/.', 14, GRAY)

# Generic connector types
conn_types = [
    ('CRM', 'Lead- & Kundenverwaltung', 'Leads, Kontakte, Angebote\nimportieren und synchronisieren', PRIMARY),
    ('Buchhaltung', 'Finanz-Systeme', 'Rechnungen, Zahlungen, Kontakte\nabgleichen und importieren', GREEN),
    ('Telefonie', 'Call-Systeme', 'Anruf-Aufzeichnungen, CDR\nautomatisch einlesen', ACCENT),
    ('Kalender', 'Termin-Systeme', 'Termine synchronisieren\nzwischen Teams und Kunden', PURPLE),
    ('E-Mail', 'Mail-Systeme', 'Automatische Benachrichtigungen\nund Follow-up Mails', RED),
    ('Custom', 'Eigene APIs', 'Beliebige REST/Webhook-\nIntegrationen moglich', GRAY),
]

for i, (name, category, desc, color) in enumerate(conn_types):
    col = i % 3
    row = i // 3
    x = Inches(0.5 + col * 4.1)
    y = Inches(1.6 + row * 2.5)
    add_box(slide, x, y, Inches(3.7), Inches(2.1), WHITE, BORDER, 0.03)
    add_box(slide, x, y, Inches(3.7), Inches(0.08), color)
    add_text(slide, x + Inches(0.3), y + Inches(0.25), Inches(2.8), Inches(0.4),
             name, 17, DARK, bold=True)
    add_text(slide, x + Inches(0.3), y + Inches(0.6), Inches(2.8), Inches(0.3),
             category, 11, color, bold=True)
    add_text(slide, x + Inches(0.3), y + Inches(1.05), Inches(3.1), Inches(0.9),
             desc, 12, GRAY)

add_box(slide, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.8), RGBColor(0xFE, 0xF3, 0xC7), BORDER, 0.02)
add_text(slide, Inches(0.8), Inches(6.4), Inches(11.5), Inches(0.6),
         'Pro Firma: company_connectors Tabelle bestimmt welche Systeme angebunden sind + API Keys/OAuth Tokens.', 14, RGBColor(0x92, 0x40, 0x0E), bold=True)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 8: Bot-Architektur
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_box(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), ACCENT)

add_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.5),
         'Bots: Pro Firma konfigurierbar', 28, WHITE, bold=True)

add_text(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.5),
         'Jede Firma entscheidet welche Bots aktiv sind. Gleiche Architektur wie Connectors.', 15, GRAY)

bots = [
    ('Setter-Bot', 'Vertrieb automatisieren', [
        'Outbound-Calls planen',
        'Lead-Qualifizierung (KI)',
        'Termine buchen',
        'Gesprache bewerten',
    ], PRIMARY),
    ('Berater-Bot', 'Angebots-Assistent', [
        'Angebote vorbereiten',
        'Follow-up Erinnerungen',
        'Pipeline optimieren',
        'Nachfass-Aktionen',
    ], GREEN),
    ('Finance-Bot', 'Finanz-Automatisierung', [
        'Rechnungen abgleichen',
        'Mahnwesen automatisieren',
        'Cashflow prognostizieren',
        'Anomalien erkennen',
    ], PURPLE),
]

for i, (name, subtitle, features, color) in enumerate(bots):
    x = Inches(0.5 + i * 4.2)
    add_box(slide, x, Inches(1.8), Inches(3.8), Inches(4.0), RGBColor(0x1F, 0x29, 0x37), radius=0.03)
    add_box(slide, x, Inches(1.8), Inches(3.8), Inches(0.08), color)
    add_text(slide, x + Inches(0.3), Inches(2.1), Inches(3.2), Inches(0.5), name, 20, WHITE, bold=True)
    add_text(slide, x + Inches(0.3), Inches(2.6), Inches(3.2), Inches(0.4), subtitle, 12, color)
    for j, feat in enumerate(features):
        add_text(slide, x + Inches(0.3), Inches(3.2 + j * 0.5), Inches(3.2), Inches(0.4),
                 f'•  {feat}', 12, RGBColor(0x9C, 0xA3, 0xAF))

add_box(slide, Inches(0.5), Inches(6.0), Inches(12.3), Inches(1.1), RGBColor(0x1F, 0x29, 0x37), ACCENT, 0.02)
add_text(slide, Inches(0.8), Inches(6.1), Inches(11.5), Inches(0.9),
         'Firma A: nur Setter-Bot  |  Firma B: alle 3 Bots  |  Firma C: eigener Custom-Bot\nBots sind Worker in apps/api/ — sie nutzen dieselben DataAccess APIs und Connector-Daten wie das Frontend.', 13, RGBColor(0x9C, 0xA3, 0xAF))

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 9: Datenfluss
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_box(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), PRIMARY)

add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.5),
         'Datenfluss: Request → Render', 28, DARK, bold=True)

steps = [
    ('1  Middleware', 'middleware.ts', 'Subdomain extrahieren → Firma aus DB laden → Branding als HTTP-Header setzen', PRIMARY),
    ('2  Layout', 'layout.tsx', 'Session laden (Profil + Rollen + Permissions) → Sidebar rendern (gefiltert nach Rechten)', RGBColor(0x04, 0x7B, 0x57)),
    ('3  Page', 'z.B. /leads', 'Permission prufen → Daten laden uber DataAccess (gleiche API fur alle Firmen)', PURPLE),
    ('4  Client', 'Interactive UI', 'Vorbereitete Daten anzeigen — kein direkter DB-Zugriff im Browser', ACCENT),
]

for i, (title, file, desc, color) in enumerate(steps):
    y = Inches(1.3 + i * 1.4)
    add_box(slide, Inches(0.8), y, Inches(0.5), Inches(0.5), color, radius=0.5)
    add_text(slide, Inches(0.8), y + Inches(0.05), Inches(0.5), Inches(0.4),
             str(i + 1), 16, WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(1.6), y, Inches(3.0), Inches(0.4), title, 16, DARK, bold=True)
    add_text(slide, Inches(1.6), y + Inches(0.4), Inches(2.5), Inches(0.3), file, 11, PRIMARY, font_name='Consolas')
    add_text(slide, Inches(4.5), y + Inches(0.05), Inches(5.5), Inches(0.9), desc, 13, GRAY)
    if i < len(steps) - 1:
        add_box(slide, Inches(1.02), y + Inches(0.55), Inches(0.06), Inches(0.85), BORDER)

add_box(slide, Inches(8.8), Inches(1.3), Inches(4.0), Inches(5.3), LIGHT_BG, BORDER, 0.02)
add_text(slide, Inches(9.1), Inches(1.45), Inches(3.5), Inches(0.4), 'DataAccess Pattern', 15, DARK, bold=True)

da_lines = [
    ('getDataAccess()', True, PRIMARY),
    ('', False, GRAY),
    ('Einheitliche API:', True, DARK),
    ('  db.offers.count(firmaId)', False, DARK),
    ('  db.leads.findAll(firmaId)', False, DARK),
    ('  db.kpis.findLatest()', False, DARK),
    ('  db.connectors.findAll()', False, DARK),
    ('', False, GRAY),
    ('2 Implementierungen:', True, DARK),
    ('  SupabaseDataAccess', False, GREEN),
    ('    → echte DB (Produktion)', False, GRAY),
    ('  MockDataAccess', False, ACCENT),
    ('    → Fake-Daten (Entwicklung)', False, GRAY),
    ('', False, GRAY),
    ('Frontend + Bots nutzen', True, RED),
    ('dieselbe API!', True, RED),
]
for i, (text, bold, color) in enumerate(da_lines):
    y = Inches(1.95 + i * 0.33)
    is_code = '(' in text or '.' in text or text.startswith('  ')
    add_text(slide, Inches(9.1), y, Inches(3.5), Inches(0.3),
             text, 10 if not bold else 11, color, bold=bold,
             font_name='Consolas' if is_code else 'Arial')

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 10: Monorepo
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_box(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), PRIMARY)

add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.5),
         'Monorepo Struktur', 28, DARK, bold=True)

tree = [
    ('enura-platform/', True),
    ('  apps/', True),
    ('    web/              Frontend (Next.js)', True),
    ('      src/middleware.ts    Firma erkennen + Branding', False),
    ('      src/lib/session.ts   Auth + Permissions', False),
    ('      src/app/(dashboard)/ Alle Seiten', False),
    ('      src/components/      UI Components', False),
    ('    api/              Backend Workers', True),
    ('      src/workers/connectors/  Connector-Worker', False),
    ('      src/workers/bots/        Bot-Worker (next)', False),
    ('  packages/', True),
    ('    types/            Shared Types', True),
    ('      src/domain.ts        Navigation + Permissions', False),
    ('      src/database.ts      DB Row Types', False),
    ('      src/data-access.ts   DataAccess Interface', False),
    ('    ui/               Shared Components', False),
    ('  supabase/', True),
    ('    migrations/        DB Migrationen', False),
]

add_box(slide, Inches(0.5), Inches(1.1), Inches(7.0), Inches(5.8), RGBColor(0x1E, 0x1E, 0x2E), radius=0.02)

for i, (line, is_folder) in enumerate(tree):
    y = Inches(1.25 + i * 0.3)
    color = ACCENT if is_folder else RGBColor(0xA5, 0xB4, 0xFC)
    add_text(slide, Inches(0.7), y, Inches(6.5), Inches(0.3),
             line, 11, color, bold=is_folder, font_name='Consolas')

# Right: wo kommen bots hin
add_box(slide, Inches(7.8), Inches(1.1), Inches(5.0), Inches(5.8), LIGHT_BG, BORDER, 0.02)
add_text(slide, Inches(8.1), Inches(1.25), Inches(4.5), Inches(0.4), 'Wo kommen neue Features hin?', 16, DARK, bold=True)

where = [
    ('Neuer Connector:', True, DARK),
    ('  apps/api/src/workers/connectors/', False, PRIMARY),
    ('  + company_connectors Eintrag in DB', False, GRAY),
    ('', False, GRAY),
    ('Neuer Bot:', True, DARK),
    ('  apps/api/src/workers/bots/', False, PRIMARY),
    ('  + company_bots Eintrag in DB', False, GRAY),
    ('  + optional: neue Sidebar-Seite', False, GRAY),
    ('', False, GRAY),
    ('Neues Modul / Seite:', True, DARK),
    ('  apps/web/src/app/(dashboard)/xyz/', False, PRIMARY),
    ('  + Permission Key in domain.ts', False, GRAY),
    ('  + navigationSections erweitern', False, GRAY),
    ('', False, GRAY),
    ('Neuer DB-Type:', True, DARK),
    ('  packages/types/src/database.ts', False, PRIMARY),
    ('  + DataAccess Interface erweitern', False, GRAY),
]

for i, (text, bold, color) in enumerate(where):
    y = Inches(1.8 + i * 0.3)
    add_text(slide, Inches(8.1), y, Inches(4.5), Inches(0.3),
             text, 11, color, bold=bold, font_name='Consolas' if text.startswith('  ') else 'Arial')

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 11: Zusammenfassung
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_box(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), PRIMARY)

add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.5),
         'Zusammenfassung', 28, WHITE, bold=True)

items = [
    ('1 Codebase, N Firmen', 'Gleicher Code — Branding, Module, Connectors, Bots alles individuell pro Firma'),
    ('Subdomain = Firma', '{firma}.enura-group.com → Middleware ladt automatisch die richtige Konfiguration'),
    ('Permission-gesteuertes UI', 'Jede Firma sieht nur die Module die freigeschaltet sind'),
    ('Pluggable Connectors', 'Beliebige externe Systeme pro Firma anbindbar — CRM, Buchhaltung, Telefonie, ...'),
    ('Bot-ready Architektur', 'Bots nutzen dieselbe API wie das Frontend — einfach als Worker anhangen'),
    ('Branchenunabhangig', 'Kein Firmenname und keine Branche im Code — reine Konfiguration'),
]

for i, (title, desc) in enumerate(items):
    y = Inches(1.2 + i * 0.95)
    add_box(slide, Inches(0.8), y, Inches(0.08), Inches(0.55), PRIMARY)
    add_text(slide, Inches(1.2), y, Inches(5.0), Inches(0.4), title, 17, WHITE, bold=True)
    add_text(slide, Inches(1.2), y + Inches(0.4), Inches(10.5), Inches(0.4), desc, 13, GRAY)

add_box(slide, Inches(0.5), Inches(6.7), Inches(8.0), Inches(0.5), PRIMARY, radius=0.03)
add_text(slide, Inches(0.5), Inches(6.73), Inches(8.0), Inches(0.4),
         '  Next Step: Bots bauen', 16, WHITE, bold=True)

# Save
output = r'C:\Users\info\enura-platform\Enura_Platform_Architecture.pptx'
prs.save(output)
print(f'Saved: {output}')
